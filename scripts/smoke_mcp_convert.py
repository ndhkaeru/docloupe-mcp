#!/usr/bin/env python3
import argparse
import asyncio
import json
import tempfile
import zipfile
from pathlib import Path

from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client


def write_pdf(path: Path, marker: str) -> None:
    content = f"BT /F1 18 Tf 72 720 Td ({marker}) Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        f"<< /Length {len(content)} >>\nstream\n{content}\nendstream".encode(),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    data = bytearray(b"%PDF-1.4\n")
    offsets = []
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(data))
        data.extend(f"{index} 0 obj\n".encode())
        data.extend(obj)
        data.extend(b"\nendobj\n")
    xref_offset = len(data)
    data.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    data.extend(b"0000000000 65535 f \n")
    for offset in offsets:
        data.extend(f"{offset:010d} 00000 n \n".encode())
    data.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n".encode()
    )
    path.write_bytes(data)


def write_docx(path: Path, marker: str) -> None:
    content_types = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>"""
    relationships = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""
    document = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body><w:p><w:r><w:t>{marker}</w:t></w:r></w:p><w:sectPr/></w:body>
</w:document>"""
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", relationships)
        archive.writestr("word/document.xml", document)


def create_sample(kind: str, root: Path) -> tuple[Path, str]:
    marker = f"DocLoupe {kind.upper()} smoke"
    if kind == "text":
        path = root / "sample.txt"
        path.write_text(marker, encoding="utf-8")
    elif kind == "csv":
        path = root / "sample.csv"
        path.write_text(f"name,value\n{marker},1\n", encoding="utf-8")
    elif kind == "json":
        path = root / "sample.json"
        path.write_text(json.dumps({"message": marker}), encoding="utf-8")
    elif kind == "html":
        path = root / "sample.html"
        path.write_text(f"<h1>{marker}</h1>", encoding="utf-8")
    elif kind == "excel":
        from openpyxl import Workbook

        path = root / "sample.xlsx"
        workbook = Workbook()
        workbook.active.append([marker])
        workbook.save(path)
    elif kind == "pptx":
        from pptx import Presentation

        path = root / "sample.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = marker
        presentation.save(path)
    elif kind == "docx":
        path = root / "sample.docx"
        write_docx(path, marker)
    elif kind == "pdf":
        path = root / "sample.pdf"
        write_pdf(path, marker)
    else:
        raise ValueError(f"Unsupported sample kind: {kind}")
    return path, marker


async def smoke(binary: Path, server_args: list[str], kind: str, timeout: float) -> None:
    with tempfile.TemporaryDirectory(prefix="docloupe-mcp-smoke-") as temp_dir:
        root = Path(temp_dir)
        sample, marker = create_sample(kind, root)
        parameters = StdioServerParameters(command=str(binary), args=server_args)
        stderr_path = root / "server-stderr.log"
        async with asyncio.timeout(timeout):
            with stderr_path.open("w", encoding="utf-8", errors="replace") as errlog:
                async with stdio_client(parameters, errlog=errlog) as (read_stream, write_stream):
                    async with ClientSession(read_stream, write_stream) as session:
                        await session.initialize()
                        tools = await session.list_tools()
                        if not any(tool.name == "convert_to_markdown" for tool in tools.tools):
                            raise RuntimeError(f"{binary.name} does not expose convert_to_markdown")
                        result = await session.call_tool("convert_to_markdown", {"file_path": str(sample)})
        stderr = stderr_path.read_text(encoding="utf-8", errors="replace") if stderr_path.exists() else ""
        if result.isError:
            raise RuntimeError(f"{binary.name} convert_to_markdown failed: {result.content}\n{stderr}")
        output = "\n".join(getattr(item, "text", "") for item in result.content)
        if marker not in output:
            raise RuntimeError(f"{binary.name} output did not contain {marker!r}: {output[:500]}\n{stderr}")
        print(f"PASS {binary.name}: convert_to_markdown ({kind})")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--binary", type=Path, required=True)
    parser.add_argument("--arg", action="append", default=[], dest="server_args")
    parser.add_argument(
        "--kind",
        choices=("excel", "pdf", "docx", "pptx", "csv", "html", "text", "json"),
        required=True,
    )
    parser.add_argument("--timeout", type=float, default=60)
    args = parser.parse_args()
    binary = args.binary.resolve()
    if not binary.is_file():
        parser.error(f"Binary not found: {binary}")
    asyncio.run(smoke(binary, args.server_args, args.kind, args.timeout))


if __name__ == "__main__":
    main()
